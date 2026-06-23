# -*- coding: utf-8 -*-
"""Erzeugt die Geschäftsführungs-Präsentation für NMGone (V2.0)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette (NMGone-Hausfarben) ----
NAVY      = RGBColor(0x0B, 0x4A, 0x86)
DARKNAVY  = RGBColor(0x07, 0x2E, 0x52)
MIDBLUE   = RGBColor(0x38, 0x67, 0xB7)
ICE       = RGBColor(0xE8, 0xF1, 0xFB)
LIGHT     = RGBColor(0xF5, 0xF7, 0xFB)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1B, 0x26, 0x33)
GREY      = RGBColor(0x5A, 0x64, 0x73)
GREEN     = RGBColor(0x11, 0x82, 0x3B)
AMBER     = RGBColor(0x8B, 0x5A, 0x00)
PURPLE    = RGBColor(0x6B, 0x4F, 0xB3)
TEAL      = RGBColor(0x0B, 0x6E, 0x6E)
LINE_GREY = RGBColor(0xD8, 0xDF, 0xEA)

HEAD = "Trebuchet MS"
BODY = "Calibri"

ROOT   = r"C:\nmg_analyse_3_1_0 - Kopie"
ASSETS = os.path.join(ROOT, "assets")
SHOTS  = os.path.join(ROOT, "praesentation_assets")
LOGO_PROD = os.path.join(ASSETS, "NMGone.png")     # 1649x954
LOGO_CORP = os.path.join(ASSETS, "nmg_logo.png")   # 2048x439  (NMG Pharma)
SHOT_DASH  = os.path.join(SHOTS, "shot_dashboard.png")  # 1056x679
SHOT_APPS  = os.path.join(SHOTS, "shot_apps.png")       # 1056x679
SHOT_KASSE = os.path.join(SHOTS, "shot_kasse.png")      # 1056x699

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.find(qn('a:effectLst'))   # inherit=False legt bereits ein leeres effectLst an
        if ef is None:
            ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '38100', 'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '1B2633'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '18000'})
        clr.append(alpha); sh.append(clr); ef.append(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb


def R(txt, size, color, bold=False, italic=False, font=BODY):
    return (txt, size, color, bold, italic, font)


def picture(s, path, x, y, w=None, h=None):
    kw = {}
    if w is not None: kw['width'] = Inches(w)
    if h is not None: kw['height'] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)


def framed_shot(s, path, x, y, w, aspect, caption=None, cap_color=GREY, cap_y=None):
    """Bild mit weißem Rahmen + Schatten; aspect = breite/höhe."""
    h = w / aspect
    pad = 0.09
    rect(s, x-pad, y-pad, w+2*pad, h+2*pad, fill=WHITE, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03, shadow=True)
    picture(s, path, x, y, w=w)
    if caption:
        cy = cap_y if cap_y is not None else (y+h+pad+0.12)
        text(s, x-pad, cy, w+2*pad, 0.7,
             [[R(caption, 12.5, cap_color)]], align=PP_ALIGN.CENTER, line_spacing=1.05)
    return h


def page_num(s, n):
    text(s, SW - 0.9, SH - 0.45, 0.6, 0.3, [[R(str(n), 10, GREY)]], align=PP_ALIGN.RIGHT)


SHOT_AR = 1056/679   # Dashboard / Apps
KASSE_AR = 1056/699  # Kasse

# =====================================================================
# 1 — TITEL
# =====================================================================
s = slide(DARKNAVY)
rect(s, 0, 0, SW, SH, fill=DARKNAVY)
# Logo-Karte (weiß) mit NMGone-Produktlogo
rect(s, 1.0, 1.05, 4.0, 2.15, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, shadow=True)
lw = 2.66; lh = lw / (1649/954)
picture(s, LOGO_PROD, 1.0 + (4.0-lw)/2, 1.05 + (2.15-lh)/2, w=lw)
# Akzentpunkte
motif = [MIDBLUE, GREEN, AMBER, PURPLE, TEAL]
for i, c in enumerate(motif):
    rect(s, 5.4 + i*0.42, 1.95, 0.28, 0.28, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
text(s, 1.0, 3.65, 11.4, 1.5, [
    [R("Die digitale Plattform für Marktanalyse,", 28, ICE, bold=True, font=HEAD)],
    [R("Rabattmanagement und Warenwirtschaft", 28, ICE, bold=True, font=HEAD)],
], space_after=4)
rect(s, 1.0, 5.35, 3.2, 0.045, fill=MIDBLUE)
text(s, 1.0, 5.65, 11.3, 1.2, [
    [R("Eigenentwicklung  ·  Stand V2.0  ·  Juni 2026", 16, RGBColor(0xB9,0xCD,0xE8), bold=True)],
    [R("Präsentation für die Geschäftsführung", 15, RGBColor(0x9F,0xB6,0xD6))],
], space_after=4)

# =====================================================================
# 2 — MANAGEMENT SUMMARY
# =====================================================================
s = slide(WHITE)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Auf einen Blick", 36, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.45, 11.5, 1.1, [
    [R("NMGone bündelt alle wiederkehrenden Aufgaben rund um Apotheken-Marktanalyse, "
       "Rabatte, Kundendaten und Verkauf in ", 17, INK),
     R("einem einzigen, selbst entwickelten Programm", 17, NAVY, bold=True),
     R(".", 17, INK)],
], line_spacing=1.15)

stats = [
    ("14", "Module unter einer\nOberfläche", NAVY),
    ("1", "Installation, eine\ngemeinsame Datenbasis", GREEN),
    ("V2.0", "Aktueller Stand –\nKasse inklusive", AMBER),
    ("0 €", "laufende Lizenz-\nkosten Dritter", PURPLE),
]
cx, cw, gap = 0.9, 2.65, 0.31
for i, (big, lab, col) in enumerate(stats):
    x = cx + i*(cw+gap)
    rect(s, x, 2.95, cw, 2.5, fill=LIGHT, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
    rect(s, x, 2.95, cw, 0.13, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
    big_sz = 44 if len(big) <= 3 else (32 if len(big) <= 5 else 26)
    text(s, x+0.15, 3.35, cw-0.3, 1.0, [[R(big, big_sz, col, bold=True, font=HEAD)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.2, 4.45, cw-0.4, 0.95, [[R(l, 13.5, GREY)] for l in lab.split("\n")],
         align=PP_ALIGN.CENTER, space_after=1)
text(s, 0.9, 5.95, 11.5, 0.8, [
    [R("Kernnutzen:  ", 15, NAVY, bold=True),
     R("maßgeschneidert auf unsere Prozesse, vollständig protokolliert (wer / was / wann), "
       "Daten bleiben im Haus.", 15, INK)],
], line_spacing=1.1)
page_num(s, 2)

# =====================================================================
# 3 — AUSGANGSLAGE / PROBLEM
# =====================================================================
s = slide(WHITE)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Warum NMGone?", 36, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.45, 11.5, 0.6, [[R("Die Ausgangslage – verteilte Werkzeuge, manueller Aufwand", 17, GREY, italic=True)]])

left = [
    ("Verteilte Excel-Dateien", "Auswertungen, Rabatte und Kundenlisten lagen in vielen einzelnen Dateien – schwer aktuell zu halten."),
    ("Manuelle Auswertung", "Marktanalysen wurden von Hand zusammengestellt – zeitaufwendig und fehleranfällig."),
    ("Keine Nachvollziehbarkeit", "Änderungen ließen sich kaum zurückverfolgen – wer hat wann was geändert?"),
    ("Kein durchgängiger Verkaufsweg", "Wareneingang, Lager und Verkauf waren nicht verbunden."),
]
y0 = 2.4
for i, (h, d) in enumerate(left):
    y = y0 + i*1.12
    rect(s, 0.9, y, 0.5, 0.5, fill=ICE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, 0.9, y, 0.5, 0.5, [[R(str(i+1), 20, NAVY, bold=True, font=HEAD)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.6, y-0.05, 5.4, 1.05, [
        [R(h, 16, INK, bold=True)], [R(d, 13, GREY)],
    ], space_after=2, line_spacing=1.05)

rect(s, 7.55, 2.4, 4.85, 4.5, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04, shadow=True)
text(s, 7.95, 2.8, 4.05, 0.6, [[R("Die Antwort", 22, WHITE, bold=True, font=HEAD)]])
text(s, 7.95, 3.5, 4.05, 3.2, [
    [R("Ein Programm statt vieler Insellösungen.", 15, ICE, bold=True)],
    [R("", 6, ICE)],
    [R("Eine gemeinsame Datenbasis für alle Module.", 14.5, WHITE)],
    [R("", 6, ICE)],
    [R("Jede Änderung wird automatisch protokolliert.", 14.5, WHITE)],
    [R("", 6, ICE)],
    [R("Von der Marktanalyse bis zum fertigen Verkauf – durchgängig.", 14.5, WHITE)],
], space_after=3, line_spacing=1.12)
page_num(s, 3)

# =====================================================================
# 4 — DIE PLATTFORM: MODULE
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.5, 11.5, 0.85, [[R("Eine Plattform – viele Module", 34, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.32, 11.5, 0.5, [[R("Alle Werkzeuge erreichbar über ein gemeinsames Dashboard", 16, GREY, italic=True)]])

modules = [
    ("Bedarfsanalyse", "PK-/ZW-Marktanalysen starten", NAVY),
    ("Produktanalyse", "Produktchancen erkennen", GREEN),
    ("NMG-Rabatte", "Rabatte, Statistik, Verlauf", PURPLE),
    ("PK-Konditionen", "Kundenindividuelle Rabatte", NAVY),
    ("Kasse", "Verkauf + Wareneingang", AMBER),
    ("Kunden", "Kundenstamm & Historie", NAVY),
    ("Schulbank", "Lernvorschläge pflegen", GREEN),
    ("Vergleichs-Suche", "PZN/Artikel schnell finden", TEAL),
    ("Globale Suche", "Kunde, Analyse, Artikel", NAVY),
    ("Abweichungsanalyse", "Manuell vs. Programm", AMBER),
    ("Daten aktualisieren", "Stammdaten & Importe", GREY),
    ("Roadmap & ToDo", "Aufgaben im Blick", PURPLE),
]
cols = 4
gw, gh = 2.6, 1.42
gx, gy = 0.9, 2.05
hgap = (SW - 2*gx - cols*gw) / (cols-1)
vgap = 0.22
for idx, (h, d, col) in enumerate(modules):
    r, c = divmod(idx, cols)
    x = gx + c*(gw+hgap); y = gy + r*(gh+vgap)
    rect(s, x, y, gw, gh, fill=WHITE, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08, shadow=True)
    rect(s, x+0.18, y+0.2, 0.16, gh-0.4, fill=col)
    text(s, x+0.46, y+0.18, gw-0.6, 0.5, [[R(h, 14.5, INK, bold=True)]])
    text(s, x+0.46, y+0.7, gw-0.6, 0.6, [[R(d, 12, GREY)]])
page_num(s, 4)

# =====================================================================
# 5 — SO SIEHT NMGone AUS (Screenshots)
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.5, 11.5, 0.85, [[R("So sieht NMGone aus", 34, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.32, 11.5, 0.5, [[R("Klare, einheitliche Oberfläche – drei Ansichten aus dem laufenden Programm", 16, GREY, italic=True)]])
iw = 3.8; ty5 = 2.3; cy5 = 5.05
framed_shot(s, SHOT_DASH,  0.49,  ty5, iw, SHOT_AR,
            "Startseite – Schnellzugriff & globale Suche", cap_y=cy5)
framed_shot(s, SHOT_KASSE, 4.767, ty5, iw, KASSE_AR,
            "Kasse – Verkauf, Wareneingang & Vorbestellungen", cap_y=cy5)
framed_shot(s, SHOT_APPS,  9.044, ty5, iw, SHOT_AR,
            "Apps – alle Module mit einem Klick", cap_y=cy5)
page_num(s, 5)

# =====================================================================
# 6 — KERNMODULE (Analyse / Rabatte / Kunden)
# =====================================================================
s = slide(WHITE)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Das Herzstück: Analyse & Rabatte", 34, NAVY, bold=True, font=HEAD)]])

cards = [
    ("Marktanalyse", NAVY, [
        "PK- und ZF-Auswertungen auf Knopfdruck",
        "Produktchancen & Abweichungs-Analyse",
        "Ergebnisse speichern und wiederfinden",
    ]),
    ("Rabatt-Management", PURPLE, [
        "NMG-Rabatte zentral pflegen – mit Verlauf",
        "Kundenindividuelle PK-Konditionen",
        "Rabatt-Kaskade automatisch im Verkauf",
    ]),
    ("Kunden & Wissen", GREEN, [
        "Kundenstamm mit kompletter Historie",
        "Schulbank: lernende Artikel-Zuordnung",
        "Globale Suche über alle Datenbestände",
    ]),
]
cw = 3.74; cx = 0.9; gap = 0.27
for i, (h, col, items) in enumerate(cards):
    x = cx + i*(cw+gap)
    rect(s, x, 1.75, cw, 4.9, fill=LIGHT, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
    rect(s, x, 1.75, cw, 0.85, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, x, 2.2, cw, 0.4, fill=col)
    text(s, x+0.3, 1.75, cw-0.6, 0.85, [[R(h, 19, WHITE, bold=True, font=HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    tb = s.shapes.add_textbox(Inches(x+0.3), Inches(2.95), Inches(cw-0.6), Inches(3.4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10); p.line_spacing = 1.08
        rr = p.add_run(); rr.text = "—  "; rr.font.bold=True
        rr.font.size=Pt(14); rr.font.color.rgb=col; rr.font.name=BODY
        r2 = p.add_run(); r2.text = it
        r2.font.size=Pt(14); r2.font.color.rgb=INK; r2.font.name=BODY
page_num(s, 6)

# =====================================================================
# 7 — DIE KASSE (mit Screenshot)
# =====================================================================
s = slide(DARKNAVY)
rect(s, 0, 0, SW, SH, fill=DARKNAVY)
rect(s, 0.9, 0.55, 2.05, 0.48, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, 0.9, 0.55, 2.05, 0.48, [[R("NEU: DIE KASSE", 11.5, WHITE, bold=True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 1.15, 6.2, 1.3, [[R("Vom Wareneingang bis zum Verkauf", 27, WHITE, bold=True, font=HEAD)]],
     line_spacing=1.05)

flow = [
    ("Wareneingang", "Artikel kommen rein: PZN, Charge, Verfall, Menge", GREEN),
    ("Lagerbestand", "Bestand pro Charge & Verfall wird mitgeführt", MIDBLUE),
    ("Verkauf", "An Apotheken verkaufen, Lager bucht automatisch ab", AMBER),
]
fy = 2.55
for i, (h, d, col) in enumerate(flow):
    y = fy + i*1.0
    rect(s, 0.9, y, 0.62, 0.62, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    text(s, 0.9, y, 0.62, 0.62, [[R(str(i+1), 24, WHITE, bold=True, font=HEAD)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.68, y-0.02, 4.7, 0.95, [
        [R(h, 17, WHITE, bold=True, font=HEAD)], [R(d, 13, ICE)],
    ], space_after=2, line_spacing=1.05)

text(s, 0.9, 5.75, 5.7, 1.5, [
    [R("Mehr als kassieren:  ", 13.5, AMBER, bold=True),
     R("Vorbestellungen & Disposition, Auftragsbestätigung per Druck oder Outlook-Mail, "
       "automatische Rabatte, Bestandskorrektur mit Pflicht-Begründung – alles protokolliert.",
       13.5, ICE)],
], line_spacing=1.16)

# Kasse-Screenshot rechts
sw_w = 5.85
framed_shot(s, SHOT_KASSE, 6.85, 1.15, sw_w, KASSE_AR,
            "Verkaufsmaske: Kunde, NMG-Artikel, Charge/Verfall, Rabatt – Beleg per Klick",
            cap_color=RGBColor(0xB9,0xCD,0xE8))
page_num(s, 7)

# =====================================================================
# 8 — NUTZEN FÜR DAS UNTERNEHMEN
# =====================================================================
s = slide(WHITE)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Was es dem Unternehmen bringt", 34, NAVY, bold=True, font=HEAD)]])

benefits = [
    ("Zeit gespart", "Auswertungen und Verkaufsbelege entstehen in Minuten statt in Stunden Handarbeit.", NAVY),
    ("Nachvollziehbar", "Jede Änderung mit Mitarbeiter, Zeitpunkt und Vorher/Nachher revisionssicher protokolliert.", GREEN),
    ("Weniger Fehler", "Automatische Rabatt-Logik und Bestandsführung statt fehleranfälliger Handeingaben.", PURPLE),
    ("Daten im Haus", "Lokale Datenbank – keine Cloud, keine sensiblen Apotheken-Daten bei Dritten.", AMBER),
    ("Keine Fremdkosten", "Eigenentwicklung – keine laufenden Lizenzgebühren für externe Software.", TEAL),
    ("Maßgeschneidert", "Exakt auf unsere Prozesse zugeschnitten und jederzeit erweiterbar.", MIDBLUE),
]
cols = 3; cw = 3.6; ch = 2.15; gx = 0.9; gy = 1.75
hgap = (SW - 2*gx - cols*cw)/(cols-1); vgap = 0.3
for idx, (h, d, col) in enumerate(benefits):
    r, c = divmod(idx, cols)
    x = gx + c*(cw+hgap); y = gy + r*(ch+vgap)
    rect(s, x, y, cw, ch, fill=WHITE, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
    rect(s, x+0.28, y+0.28, 0.55, 0.55, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, x+0.28, y+0.28, 0.55, 0.55, [[R("✓", 22, WHITE, bold=True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+1.0, y+0.32, cw-1.25, 0.55, [[R(h, 17, INK, bold=True, font=HEAD)]])
    text(s, x+0.3, y+1.05, cw-0.6, 1.0, [[R(d, 13, GREY)]], line_spacing=1.1)
page_num(s, 8)

# =====================================================================
# 9 — TECHNIK & SICHERHEIT
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Technik & Sicherheit – kurz erklärt", 32, NAVY, bold=True, font=HEAD)]])

facts = [
    ("Windows-Programm", "Läuft lokal auf den Arbeitsplätzen, Installation per Setup – kein Server nötig."),
    ("Gemeinsame Datenbank", "Alle Module teilen eine zentrale Datenbasis (SQLite, mehrplatzfähig)."),
    ("Automatische Updates", "Versionsschema mit Meilensteinen & Service-Packs; Updates per Installer."),
    ("Datensicherung", "Integrierte Backups; sensible Daten verlassen das Haus nicht."),
    ("Änderungsprotokoll", "Wer-was-wann über alle Module – Grundlage für Auswertung & Revision."),
    ("Eigene Codebasis", "Vollständig im Haus entwickelt – unabhängig von Fremdanbietern."),
]
for idx, (h, d) in enumerate(facts):
    r, c = divmod(idx, 2)
    x = 0.9 + c*6.05; y = 1.85 + r*1.55
    rect(s, x, y, 5.7, 1.32, fill=WHITE, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
    rect(s, x, y, 0.16, 1.32, fill=NAVY)
    text(s, x+0.45, y+0.2, 5.0, 0.5, [[R(h, 16, NAVY, bold=True, font=HEAD)]])
    text(s, x+0.45, y+0.68, 5.0, 0.6, [[R(d, 12.5, GREY)]], line_spacing=1.05)
page_num(s, 9)

# =====================================================================
# 10 — ROADMAP / AUSBLICK
# =====================================================================
s = slide(WHITE)
text(s, 0.9, 0.55, 11.5, 0.9, [[R("Ausblick", 36, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.4, 11.5, 0.5, [[R("Die Plattform wächst weiter – die nächsten Schritte", 16, GREY, italic=True)]])

bw, half = 2.75, 1.375
c0, c3 = 1.925, 11.408
ty = 2.95
rect(s, c0-0.1, ty, (c3 - c0) + 0.2, 0.06, fill=LINE_GREY)
steps = [
    ("JETZT", "V2.0 im Einsatz", "Bedarfsanalyse, Mitarbeiter,\nHilfe & Report neu dabei", GREEN),
    ("IN ARBEIT", "Organigramm-Grafik", "Mitarbeiter-Struktur\nvisuell darstellen", NAVY),
    ("GEPLANT", "Biosimilar-Wissen", "Wirkstoff-Zuordnung in\nder Bedarfsanalyse", PURPLE),
    ("DAS GROSSE ZIEL", "Cloud-basiertes Arbeiten", "Zentrale Daten, mehrplatz-\nfähig & weltweit  →  Folie 11", TEAL),
]
for i, (when, h, d, col) in enumerate(steps):
    cxp = c0 + i*((c3 - c0)/3)
    text(s, cxp-half, ty-0.5, bw, 0.32, [[R(when, 11.5, col, bold=True)]], align=PP_ALIGN.CENTER)
    rect(s, cxp-0.15, ty-0.12, 0.3, 0.3, fill=col, shape=MSO_SHAPE.OVAL)
    by = ty + 0.55
    rect(s, cxp-half, by, bw, 2.55, fill=LIGHT, line=LINE_GREY, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, shadow=True)
    rect(s, cxp-half, by, bw, 0.12, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
    text(s, cxp-half+0.15, by+0.42, bw-0.3, 0.7, [[R(h, 15.5, INK, bold=True, font=HEAD)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxp-half+0.18, by+1.35, bw-0.36, 1.0,
         [[R(l, 12.5, GREY)] for l in d.split("\n")],
         align=PP_ALIGN.CENTER, space_after=1, line_spacing=1.05)
page_num(s, 10)

# =====================================================================
# 11 — VISION: CLOUD-BASIERTES ARBEITEN
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.45, 11.6, 0.8, [[R("Cloud-basiertes Arbeiten", 32, NAVY, bold=True, font=HEAD)]])
text(s, 0.9, 1.25, 11.6, 0.5, [[R("Programm lokal, Daten zentral – neue Rabatte oder Kunden sind sofort überall sichtbar, im Büro wie unterwegs", 15, GREY, italic=True)]])

# Linke Spalte: Nutzen
cloud_benefits = [
    ("Immer der aktuelle Stand", "Neue Rabattliste oder neuer Kunde – sofort bei allen, Büro wie Kasse.", NAVY),
    ("Kasse weltweit", "Als Web-App ohne App-Store – überall sofort einsatzbereit.", GREEN),
    ("DSGVO & EU-Hosting", "Daten in EU-Region, verschlüsselt und per Login abgesichert.", TEAL),
    ("Schritt für Schritt", "Aufbau in Phasen mit geringem Risiko – Konzept & Test stehen bereits.", AMBER),
]
by = 2.05
for i, (h, d, col) in enumerate(cloud_benefits):
    y = by + i*0.95
    rect(s, 0.9, y, 0.5, 0.5, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, 0.9, y, 0.5, 0.5, [[R("✓", 19, WHITE, bold=True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.6, y-0.04, 4.95, 0.95, [[R(h, 15.5, INK, bold=True)], [R(d, 12.5, GREY)]],
         space_after=1, line_spacing=1.05)

# Rechte Spalte: Architektur-Diagramm
def diagbox(x, y, w, h, title, sub, col, fill=WHITE, tcol=INK, scol=GREY):
    rect(s, x, y, w, h, fill=fill, line=(col if fill == WHITE else None), line_w=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, shadow=True)
    if fill == WHITE:
        rect(s, x, y, 0.13, h, fill=col)
    text(s, x+0.3, y, w-0.5, h, [[R(title, 13.5, tcol, bold=True)], [R(sub, 11, scol)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=1, line_spacing=1.0)

diagbox(6.95, 1.95, 2.7, 0.9, "NMGone", "Büro · lokal", NAVY)
diagbox(9.95, 1.95, 2.7, 0.9, "Kasse", "Web-App · weltweit", GREEN)
for ax in (8.13, 11.13):
    rect(s, ax, 2.95, 0.34, 0.5, fill=MIDBLUE, shape=MSO_SHAPE.DOWN_ARROW)
diagbox(7.4, 3.55, 4.85, 0.9, "API-Backend · EU-Cloud", "einziger, abgesicherter Zugang",
        NAVY, fill=NAVY, tcol=WHITE, scol=ICE)
rect(s, 9.65, 4.55, 0.36, 0.48, fill=MIDBLUE, shape=MSO_SHAPE.DOWN_ARROW)
diagbox(7.4, 5.1, 4.85, 0.9, "Zentrale Datenbank · PostgreSQL (EU)", "die eine Quelle der Wahrheit",
        TEAL, fill=TEAL, tcol=WHITE, scol=ICE)

# Footer: Phasen + Kosten
text(s, 0.9, 6.35, 11.6, 0.9, [
    [R("Aufbau in Phasen:  ", 12.5, NAVY, bold=True),
     R("Aufräumen → Backend → Cloud-Hosting → Umstellen → Go-Live", 12.5, INK)],
    [R("Betriebskosten moderat:  ", 12.5, NAVY, bold=True),
     R("ca. 30–80 €/Monat (EU-Hosting + Datenbank) · kein App-Store nötig", 12.5, GREY)],
], space_after=3, line_spacing=1.05)
page_num(s, 11)

# =====================================================================
# 12 — ABSCHLUSS
# =====================================================================
s = slide(DARKNAVY)
rect(s, 0, 0, SW, SH, fill=DARKNAVY)
for i, c in enumerate([NAVY, MIDBLUE, GREEN, AMBER, PURPLE, TEAL]):
    rect(s, 1.0 + i*0.62, 1.4, 0.46, 0.46, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
text(s, 1.0, 2.5, 11.3, 1.6, [
    [R("Eine Plattform. Alle Prozesse.", 40, WHITE, bold=True, font=HEAD)],
    [R("Im Haus gebaut.", 40, MIDBLUE, bold=True, font=HEAD)],
], space_after=2)
text(s, 1.0, 4.55, 11.0, 1.2, [
    [R("NMGone fasst Marktanalyse, Rabatte, Kunden und Verkauf zusammen – "
       "nachvollziehbar, kostenneutral und exakt auf unsere Apotheken-Prozesse zugeschnitten.",
       17, ICE)],
], line_spacing=1.25)
rect(s, 1.0, 6.1, 3.0, 0.045, fill=MIDBLUE)
text(s, 1.0, 6.35, 8.0, 0.6, [[R("Vielen Dank – gerne beantworte ich Ihre Fragen.", 15,
     RGBColor(0xB9,0xCD,0xE8), bold=True)]])
# Firmenlogo (NMG Pharma) in weißer Chip unten rechts
rect(s, 10.25, 6.25, 2.45, 0.8, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18, shadow=True)
clw = 2.05; clh = clw / (2048/439)
picture(s, LOGO_CORP, 10.25 + (2.45-clw)/2, 6.25 + (0.8-clh)/2, w=clw)

out = os.path.join(ROOT, "NMGone_Praesentation_Geschaeftsfuehrung.pptx")
prs.save(out)
print("OK ->", out, "Slides:", len(prs.slides._sldIdLst))
